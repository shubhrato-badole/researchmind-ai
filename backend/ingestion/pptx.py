from langchain_core.documents import Document
from ingestion.utils import process_and_store
from pptx import Presentation
import tempfile
import os

def ingest_pptx(file_bytes: bytes, filename: str, user_id: int):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name

        prs = Presentation(tmp_path)
        os.unlink(tmp_path)

        # extract text from all slides
        full_text = ""
        for i, slide in enumerate(prs.slides):
            slide_text = f"\nSlide {i+1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            full_text += slide_text

        if not full_text.strip():
            return {"message": "No text found in presentation", "chunks": 0}

        docs = [Document(page_content=full_text, metadata={"source": filename})]
        return process_and_store(docs, filename, "pptx", None, user_id)

    except Exception as e:
        return {"message": f"Error: {str(e)}", "chunks": 0}